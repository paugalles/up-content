from typing import List
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class LinkedinSlideGenerator:
    def __init__(self, config: dict):
        self.w, self.h = config.get('video', {}).get('resolution', (1920, 1080))
        self.config = config
        
        self.bg_color = RGBColor(248, 250, 252)
        self.title_color = RGBColor(12, 61, 109)
        self.accent_color = RGBColor(12, 136, 235)
        self.body_color = RGBColor(15, 23, 42)
        self.font_name = 'Plus Jakarta Sans'

    @staticmethod
    def _send_shape_to_back(slide, shape):
        shape_element = shape.element
        shape_element.getparent().remove(shape_element)
        slide.shapes._spTree.insert(2, shape_element)

    def _style_pptx_background(self, slide, is_title=False):
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = self.bg_color

        if not is_title:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(1.35), Inches(9.3), Inches(3.95)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            self._send_shape_to_back(slide, card)

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.accent_color
        accent.line.fill.background()
        self._send_shape_to_back(slide, accent)

    def generate_pptx(self, scenes, output_pptx: str):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        title_slide_layout = prs.slide_layouts[0]
        bullet_slide_layout = prs.slide_layouts[1]
        
        for scene in scenes:
            if scene.is_title:
                slide = prs.slides.add_slide(title_slide_layout)
                self._style_pptx_background(slide, is_title=True)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = scene.slide_title
                title.left = Inches(0.75)
                title.top = Inches(1.5)
                title.width = Inches(8.5)
                title.height = Inches(1.6)
                
                if title.text_frame.paragraphs:
                    p = title.text_frame.paragraphs[0]
                    p.font.name = self.font_name
                    p.font.size = Pt(48)
                    p.font.color.rgb = self.title_color
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                
                if scene.slide_body:
                    subtitle.text = scene.slide_body
                    subtitle.left = Inches(1.5)
                    subtitle.top = Inches(3.2)
                    subtitle.width = Inches(7.0)
                    subtitle.height = Inches(1.9)
                    
                    if subtitle.text_frame.paragraphs:
                        p = subtitle.text_frame.paragraphs[0]
                        p.font.name = self.font_name
                        p.font.size = Pt(24)
                        p.font.color.rgb = self.accent_color
                        p.alignment = PP_ALIGN.CENTER
                else:
                    sp = subtitle.element
                    sp.getparent().remove(sp)
                    
            else:
                slide = prs.slides.add_slide(bullet_slide_layout)
                self._style_pptx_background(slide)
                title = slide.shapes.title
                body = slide.placeholders[1]
                
                title.text = scene.slide_title
                title.left = Inches(0.5)
                title.top = Inches(0.3)
                title.width = Inches(9.0)
                title.height = Inches(1.25)
                
                if title.text_frame.paragraphs:
                    p = title.text_frame.paragraphs[0]
                    p.font.name = self.font_name
                    p.font.size = Pt(32)
                    p.font.color.rgb = self.title_color
                    p.font.bold = True
                
                body.left = Inches(0.5)
                body.top = Inches(1.75)
                body.width = Inches(9.0)
                body.height = Inches(3.5)
                
                try:
                    import json
                    data = json.loads(scene.slide_body)
                except:
                    data = {"type": "bullets", "items": [{"title": "", "content": scene.slide_body}]}
                
                if data.get("type") == "chart":
                    sp = body.element
                    sp.getparent().remove(sp)
                    
                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE
                    
                    chart_data = CategoryChartData()
                    chart_data.categories = data.get("labels", [])
                    chart_data.add_series(data.get("title", "Data"), data.get("values", []))
                    
                    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
                    if data.get("chart_type") == "pie":
                        chart_type = XL_CHART_TYPE.PIE
                    elif data.get("chart_type") == "line":
                        chart_type = XL_CHART_TYPE.LINE
                        
                    slide.shapes.add_chart(
                        chart_type, Inches(1.0), Inches(1.75), Inches(8.0), Inches(3.5), chart_data
                    )
                elif data.get("type") == "process":
                    sp = body.element
                    sp.getparent().remove(sp)
                    
                    from pptx.enum.shapes import MSO_SHAPE
                    from pptx.dml.color import RGBColor
                    from pptx.enum.text import PP_ALIGN
                    
                    steps = data.get("steps", [])
                    num_steps = len(steps)
                    if num_steps > 0:
                        spacing = 0.2
                        shape_w = (8.5 - (num_steps - 1) * spacing) / num_steps
                        start_x = 0.75
                        
                        for i, step_text in enumerate(steps):
                            x = Inches(start_x + i * (shape_w + spacing))
                            y = Inches(2.5)
                            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(shape_w), Inches(1.5))
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(12, 136, 235)
                            shape.line.color.rgb = RGBColor(12, 61, 109)
                            
                            tf = shape.text_frame
                            tf.word_wrap = True
                            p = tf.paragraphs[0]
                            p.text = step_text
                            p.font.size = Pt(16)
                            p.font.color.rgb = RGBColor(255, 255, 255)
                            p.alignment = PP_ALIGN.CENTER
                            
                            if i < num_steps - 1:
                                arrow_x = x + Inches(shape_w) + Inches(spacing/2) - Inches(0.1)
                                arrow_y = y + Inches(0.6)
                                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.2), Inches(0.3))
                                arrow.fill.solid()
                                arrow.fill.fore_color.rgb = RGBColor(15, 23, 42)
                else:
                    tf = body.text_frame
                    tf.clear()
                    
                    items = data.get("items", [])
                    texts_to_add = []
                    for item in items:
                        t = item.get("title", "")
                        c = item.get("content", "")
                        texts_to_add.append(f"{t}: {c}" if t and c else f"{t}{c}")
                    
                    if not texts_to_add:
                        raw_text = scene.slide_body
                        if "•" in raw_text:
                            texts_to_add = [t.strip() for t in raw_text.split("•") if t.strip()]
                        else:
                            texts_to_add = [raw_text]
                    
                    for i, text in enumerate(texts_to_add):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                            
                        p.text = text.replace("•", "").strip()
                        p.level = 0
                        p.font.name = self.font_name
                        p.font.size = Pt(18)
                        p.font.color.rgb = self.body_color
                    p.font.color.rgb = self.body_color

        prs.save(output_pptx)
