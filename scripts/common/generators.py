import os
import re
import random
import textwrap
import glob
from dataclasses import dataclass
from typing import List, Dict, Optional
from pathlib import Path

from openai import OpenAI
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from moviepy.editor import AudioFileClip, CompositeAudioClip, CompositeVideoClip, ImageClip, concatenate_audioclips
import moviepy.audio.fx.all as afx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from .articles import Article
from .config import branding

class Scene:
    def __init__(self, spoken_text: str, slide_title: str, slide_body: str, is_title: bool = False):
        self.spoken_text = spoken_text
        self.slide_title = slide_title
        self.slide_body = slide_body
        self.is_title = is_title

# Copy NarrationGenerator logic exactly, just adapt to `Article`
class NarrationGenerator:
    def generate_scenes(self, article: Article) -> List[Scene]:
        lang = article.language
        brand = branding().watermark_text if hasattr(branding(), 'watermark_text') else branding().name
        cta_text = branding().cta_es if lang == 'es' else branding().cta_en
        
        client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
        
        if lang == 'es':
            intro = f"Hola y bienvenidos. Hoy hablaremos sobre: {article.title}. Este video llega gracias a {brand}."
            outro = f"Nota: Este video es informativo y no constituye asesoramiento legal. Gracias por ver. {cta_text}"
        elif lang == 'en':
            intro = f"Hello and welcome. Today we are talking about: {article.title}. This video is brought to you by {brand}."
            outro = f"Note: This video is for informational purposes and does not constitute legal advice. Thanks for watching. {cta_text}"
        else:
            base_intro = f"Hello and welcome. Today we are talking about: {article.title}. This video is brought to you by {brand}."
            base_outro = f"Note: This video is for informational purposes and does not constitute legal advice. Thanks for watching. {cta_text}"
            try:
                sys_msg = {"role": "system", "content": f"Translate the following text into the language '{lang}'. Output ONLY the translated text, nothing else."}
                res_intro = client.chat.completions.create(model="gpt-4o-mini", messages=[sys_msg, {"role": "user", "content": base_intro}], max_tokens=100)
                intro = res_intro.choices[0].message.content.strip()
                res_outro = client.chat.completions.create(model="gpt-4o-mini", messages=[sys_msg, {"role": "user", "content": base_outro}], max_tokens=100)
                outro = res_outro.choices[0].message.content.strip()
            except Exception:
                intro = base_intro
                outro = base_outro
            
        scenes = []
        scenes.append(Scene(spoken_text=intro, slide_title=article.title, slide_body="", is_title=True))
        
        # Split markdown text into grouped sections
        grouped_sections = []
        current_group = {"heading": "", "paragraphs": []}
        
        lines = article.text.split('\n')
        for line in lines:
            line = line.strip()
            if not line: continue
            if line.startswith('##'):
                if current_group["heading"] or current_group["paragraphs"]:
                    grouped_sections.append(current_group)
                current_group = {"heading": line.replace('##', '').strip(), "paragraphs": []}
            elif line.startswith('#'):
                continue
            else:
                current_group["paragraphs"].append(line)
                
        if current_group["heading"] or current_group["paragraphs"]:
            grouped_sections.append(current_group)
            
        for group in grouped_sections:
            heading = group["heading"]
            paragraphs_text = " ".join(group["paragraphs"])
            
            if not paragraphs_text and heading:
                scenes.append(Scene(spoken_text=heading, slide_title=heading, slide_body="", is_title=True))
                continue
                
            if not heading:
                heading = "Puntos Clave" if lang == 'es' else "Key Points"
                
            spoken_text = f"{heading}. {paragraphs_text}"
            
            try:
                sys_prompt = f"""You are a helpful assistant that summarizes text for a PowerPoint slide. 
If the text describes comparisons between options, costs/fees, wait times (e.g., in months), or any meaningful quantitative data, output a JSON object representing a chart that visualizes this data:
{{
  "type": "chart",
  "chart_type": "pie", // can be 'bar', 'pie', or 'line'
  "title": "Chart Title",
  "labels": ["Option A", "Option B"], 
  "values": [500, 800] // Must be actual numerical values that make sense to graph
}}
If the text describes a step-by-step process or a timeline, output a JSON object representing a process diagram:
{{
  "type": "process",
  "title": "Process Title",
  "steps": ["Step 1 description", "Step 2 description", "Step 3 description"] // max 4 or 5 steps to fit on a slide
}}
Otherwise, output a JSON object with bullet points:
{{
  "type": "bullets",
  "items": [
    {{"title": "Paso 1", "content": "Details..."}},
    {{"title": "Paso 2", "content": "Details..."}}
  ]
}}
Write the content in the language '{lang}'. Output ONLY valid JSON."""
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    response_format={ "type": "json_object" },
                    messages=[
                        {"role": "system", "content": sys_prompt},
                        {"role": "user", "content": f"Please summarize this section for a presentation slide:\n\n{paragraphs_text}"}
                    ],
                    max_tokens=250
                )
                bullet_points = response.choices[0].message.content.strip()
            except Exception:
                bullet_points = "• " + paragraphs_text[:100] + "..."
                
            scenes.append(Scene(spoken_text=spoken_text, slide_title=heading, slide_body=bullet_points, is_title=False))
            
        scenes.append(Scene(spoken_text=outro, slide_title=brand, slide_body=cta_text, is_title=True))
        return scenes
class EdgeTTSProvider:
    def __init__(self, config: dict):
        self.language = config.get('language', 'es')
        if self.language.startswith('es'):
            self.voice = "es-ES-AlvaroNeural"
        elif self.language.startswith('en'):
            self.voice = "en-US-AriaNeural"
        else:
            self.voice = "es-ES-AlvaroNeural"

    async def generate(self, text: str, output_audio: str):
        print(f"Calling Edge TTS API with voice {self.voice}...")
        import edge_tts
        communicate = edge_tts.Communicate(text, self.voice)
        await communicate.save(output_audio)


class SlideGenerator:
    def __init__(self, config: dict):
        self.w, self.h = config['video']['resolution']
        self.config = config
        
        # Website-inspired palette: slate canvas, blue glow, white cards.
        self.bg_color = RGBColor(248, 250, 252)
        
        # Title Color: 0C3D6D (brand-900 dark blue)
        self.title_color = RGBColor(12, 61, 109)
        
        # Subtitle / Accent: 0C88EB (brand-500 vivid blue)
        self.accent_color = RGBColor(12, 136, 235)
        
        # Body text: 0F172A (slate-900 almost black)
        self.body_color = RGBColor(15, 23, 42)
        
        self.font_name = 'Plus Jakarta Sans'

    @staticmethod
    def _send_shape_to_back(slide, shape):
        """Place a decorative shape behind the layout placeholders."""
        shape_element = shape.element
        shape_element.getparent().remove(shape_element)
        slide.shapes._spTree.insert(2, shape_element)

    def _style_pptx_background(self, slide, is_title=False):
        """Apply the same cool slate and blue-glow language as the website."""
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = self.bg_color

        if not is_title:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(1.35), Inches(9.3), Inches(3.95)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(226, 232, 240)
            self._send_shape_to_back(slide, card)

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.accent_color
        accent.line.fill.background()
        self._send_shape_to_back(slide, accent)

    def generate_pptx(self, scenes: List[Scene], output_pptx: str):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        # Ensure 16:9 widescreen layout like your template
        prs.slide_width = Inches(10) # 16:9 ratio in inches (usually 13.33 x 7.5, but your template used 10x5.625)
        prs.slide_height = Inches(5.625)
        
        # Layouts
        # 0 = Title Slide (Centered title + subtitle)
        # 1 = Title and Content (Title + Bullet points)
        title_slide_layout = prs.slide_layouts[0]
        bullet_slide_layout = prs.slide_layouts[1]
        
        for scene in scenes:
            if scene.is_title:
                slide = prs.slides.add_slide(title_slide_layout)
                self._style_pptx_background(slide, is_title=True)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                # Title styling
                title.text = scene.slide_title
                title.left = Inches(0.75)
                title.top = Inches(1.5)
                title.width = Inches(8.5)
                title.height = Inches(1.6)
                
                if title.text_frame.paragraphs:
                    p = title.text_frame.paragraphs[0]
                    p.font.name = self.font_name
                    p.font.size = Pt(48)
                    p.font.color.rgb = self.title_color
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                
                # Subtitle styling
                if scene.slide_body:
                    subtitle.text = scene.slide_body
                    subtitle.left = Inches(1.5)
                    subtitle.top = Inches(3.2)
                    subtitle.width = Inches(7.0)
                    subtitle.height = Inches(1.9)
                    
                    if subtitle.text_frame.paragraphs:
                        p = subtitle.text_frame.paragraphs[0]
                        p.font.name = self.font_name
                        p.font.size = Pt(24)
                        p.font.color.rgb = self.accent_color
                        p.alignment = PP_ALIGN.CENTER
                else:
                    sp = subtitle.element
                    sp.getparent().remove(sp)
                    
            else:
                slide = prs.slides.add_slide(bullet_slide_layout)
                self._style_pptx_background(slide)
                title = slide.shapes.title
                body = slide.placeholders[1]
                
                title.text = scene.slide_title
                title.left = Inches(0.5)
                title.top = Inches(0.3)
                title.width = Inches(9.0)
                title.height = Inches(1.25)
                
                if title.text_frame.paragraphs:
                    p = title.text_frame.paragraphs[0]
                    p.font.name = self.font_name
                    p.font.size = Pt(32)
                    p.font.color.rgb = self.title_color
                    p.font.bold = True
                
                body.left = Inches(0.5)
                body.top = Inches(1.75)
                body.width = Inches(9.0)
                body.height = Inches(3.5)
                
                try:
                    import json
                    data = json.loads(scene.slide_body)
                except:
                    data = {"type": "bullets", "items": [{"title": "", "content": scene.slide_body}]}
                
                if data.get("type") == "chart":
                    # Remove the text placeholder and add a chart
                    sp = body.element
                    sp.getparent().remove(sp)
                    
                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE
                    
                    chart_data = CategoryChartData()
                    chart_data.categories = data.get("labels", [])
                    chart_data.add_series(data.get("title", "Data"), data.get("values", []))
                    
                    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
                    if data.get("chart_type") == "pie":
                        chart_type = XL_CHART_TYPE.PIE
                    elif data.get("chart_type") == "line":
                        chart_type = XL_CHART_TYPE.LINE
                        
                    slide.shapes.add_chart(
                        chart_type, Inches(1.0), Inches(1.75), Inches(8.0), Inches(3.5), chart_data
                    )
                elif data.get("type") == "process":
                    # Remove the text placeholder and add process shapes
                    sp = body.element
                    sp.getparent().remove(sp)
                    
                    from pptx.enum.shapes import MSO_SHAPE
                    from pptx.dml.color import RGBColor
                    from pptx.enum.text import PP_ALIGN
                    
                    steps = data.get("steps", [])
                    num_steps = len(steps)
                    if num_steps > 0:
                        spacing = 0.2
                        shape_w = (8.5 - (num_steps - 1) * spacing) / num_steps
                        start_x = 0.75
                        
                        for i, step_text in enumerate(steps):
                            x = Inches(start_x + i * (shape_w + spacing))
                            y = Inches(2.5)
                            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(shape_w), Inches(1.5))
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(12, 136, 235)
                            shape.line.color.rgb = RGBColor(12, 61, 109)
                            
                            tf = shape.text_frame
                            tf.word_wrap = True
                            p = tf.paragraphs[0]
                            p.text = step_text
                            p.font.size = Pt(16)
                            p.font.color.rgb = RGBColor(255, 255, 255)
                            p.alignment = PP_ALIGN.CENTER
                            
                            # draw arrow between steps
                            if i < num_steps - 1:
                                arrow_x = x + Inches(shape_w) + Inches(spacing/2) - Inches(0.1)
                                arrow_y = y + Inches(0.6)
                                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.2), Inches(0.3))
                                arrow.fill.solid()
                                arrow.fill.fore_color.rgb = RGBColor(15, 23, 42)
                else:
                    tf = body.text_frame
                    tf.clear()
                    
                    items = data.get("items", [])
                    texts_to_add = []
                    for item in items:
                        t = item.get("title", "")
                        c = item.get("content", "")
                        texts_to_add.append(f"{t}: {c}" if t and c else f"{t}{c}")
                    
                    if not texts_to_add:
                        # Fallback parsing
                        raw_text = scene.slide_body
                        if "•" in raw_text:
                            texts_to_add = [t.strip() for t in raw_text.split("•") if t.strip()]
                        else:
                            texts_to_add = [raw_text]
                    
                    # We render this as a standard bulleted list
                    for i, text in enumerate(texts_to_add):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                            
                        p.text = text.replace("•", "").strip() # Remove any rogue bullet chars since PPTX adds them natively
                        p.level = 0
                        p.font.name = self.font_name
                        p.font.size = Pt(18)
                        p.font.color.rgb = self.body_color
                    p.font.color.rgb = self.body_color

        prs.save(output_pptx)

    def generate_image(self, scene: Scene, output_path: str):
        # We simulate the exact styling above in Pillow for MoviePy integration
        from PIL import Image, ImageDraw, ImageFilter, ImageFont
        import textwrap
        img = Image.new('RGB', (self.w, self.h), (248, 250, 252))

        # Soft, blurred glows reproduce the website hero background without
        # competing with slide content.
        glow_layer = Image.new('RGBA', (self.w, self.h), (0, 0, 0, 0))
        glow_draw = ImageDraw.Draw(glow_layer)
        glow_draw.ellipse(
            [self.w * 0.52, -self.h * 0.35, self.w * 1.12, self.h * 0.55],
            fill=(191, 219, 254, 145),
        )
        glow_draw.ellipse(
            [-self.w * 0.18, self.h * 0.70, self.w * 0.32, self.h * 1.25],
            fill=(186, 230, 253, 105),
        )
        glow_layer = glow_layer.filter(ImageFilter.GaussianBlur(radius=95))
        img = Image.alpha_composite(img.convert('RGBA'), glow_layer).convert('RGB')
        draw = ImageDraw.Draw(img)

        draw.rectangle([0, 0, self.w, 12], fill=(12, 136, 235))
        if not scene.is_title:
            draw.rounded_rectangle(
                [65, 255, self.w - 65, self.h - 55],
                radius=42,
                fill=(255, 255, 255),
                outline=(226, 232, 240),
                width=3,
            )
        
        font_path = branding().font_path or "/System/Library/Fonts/Supplemental/Arial.ttf"
        
        # Scaling PPTX 10x5.625 inches to 1920x1080 (Scaling factor = 192)
        # Font Pt to Pixels (1pt = 1/72 inch -> * 192 = 2.66 pixels)
        # However, Pillow uses TTF sizes which roughly map 1:1 to pixels. We will scale up slightly.
        title_font = ImageFont.truetype(font_path, 90) # ~ 32pt scaled
        main_title_font = ImageFont.truetype(font_path, 130) # ~ 48pt scaled
        subtitle_font = ImageFont.truetype(font_path, 65) # ~ 24pt scaled
        body_font = ImageFont.truetype(font_path, 50) # ~ 18pt scaled
        
        title_color = (12, 61, 109)
        accent_color = (12, 136, 235)
        body_color = (15, 23, 42)
        
        def draw_wrapped_text(draw, text, font, x, y, width, fill, align="left", bullet=False):
            # If the text comes with multiple bullets embedded (e.g. from a raw string not parsed as json),
            # split it by the bullet character first so we can draw them as separate indented blocks.
            items_to_draw = [text]
            if bullet and "•" in text:
                items_to_draw = [t.strip() for t in text.split("•") if t.strip()]
                
            current_y = y
            
            for item in items_to_draw:
                lines = textwrap.wrap(item, width=width)
                for i, line in enumerate(lines):
                    bbox = draw.textbbox((0,0), line, font=font)
                    line_w = bbox[2] - bbox[0]
                    line_h = bbox[3] - bbox[1]
                    draw_x = x
                    if align == "center":
                        draw_x = x - (line_w // 2)
                        
                    if bullet and i == 0:
                        # Draw actual bullet circle manually
                        r = 8
                        cy = current_y + (line_h // 2) + 5
                        cx = draw_x - 30
                        draw.ellipse([cx-r, cy-r, cx+r, cy+r], fill=fill)
                        
                    draw.text((draw_x, current_y), line, font=font, fill=fill)
                    current_y += line_h + 15
                current_y += 30 # Extra spacing between bullet items
            return current_y

        if scene.is_title:
            draw_wrapped_text(draw, scene.slide_title, main_title_font, self.w//2, self.h//2 - 150, 25, title_color, align="center")
            if scene.slide_body:
                draw_wrapped_text(draw, scene.slide_body, subtitle_font, self.w//2, self.h//2 + 50, 45, accent_color, align="center")
        else:
            # Header
            draw_wrapped_text(draw, scene.slide_title, title_font, 100, 60, 40, title_color, align="left")
            
            try:
                import json
                data = json.loads(scene.slide_body)
            except:
                data = {"type": "bullets", "items": [{"title": "", "content": scene.slide_body}]}
            
            if data.get("type") == "chart":
                # Draw a chart using matplotlib and paste it onto the slide image
                import matplotlib.pyplot as plt
                import io
                
                fig, ax = plt.subplots(figsize=(12, 5), dpi=120)
                labels = data.get("labels", [])
                values = data.get("values", [])
                
                chart_type = data.get("chart_type", "bar")
                if chart_type == "pie":
                    ax.pie(values, labels=labels, autopct='%1.1f%%', colors=['#0C88EB', '#F59E0B', '#10B981', '#14B8A6'])
                elif chart_type == "line":
                    ax.plot(labels, values, color='#0C88EB', marker='o', linewidth=3, markersize=8)
                else: # bar
                    ax.bar(labels, values, color='#0C88EB')
                    
                ax.set_title(data.get("title", ""), fontsize=20, color='#0F172A')
                fig.tight_layout()
                
                buf = io.BytesIO()
                fig.savefig(buf, format='png', transparent=True)
                buf.seek(0)
                plt.close(fig)
                
                chart_img = Image.open(buf).convert("RGBA")
                # Paste the chart image onto the main PIL image
                img.paste(chart_img, (150, 300), chart_img)
                
            elif data.get("type") == "process":
                steps = data.get("steps", [])
                num_steps = len(steps)
                if num_steps > 0:
                    import textwrap
                    
                    spacing = 50
                    start_x = 100
                    start_y = 350
                    box_h = 400
                    
                    # Compute box width so it perfectly fills the slide horizontally
                    box_w = int((self.w - (start_x * 2) - (num_steps - 1) * spacing) / num_steps)
                    
                    step_font = ImageFont.truetype(branding().font_path or "/System/Library/Fonts/Supplemental/Arial.ttf", 40)
                    
                    for i, step_text in enumerate(steps):
                        x1 = start_x + i * (box_w + spacing)
                        y1 = start_y
                        x2 = x1 + box_w
                        y2 = y1 + box_h
                        
                        # Draw rounded rectangle (available in recent Pillow versions)
                        draw.rounded_rectangle([x1, y1, x2, y2], radius=25, fill=(12, 136, 235), outline=(12, 61, 109), width=5)
                        
                        # Calculate text wrap based on box width
                        # average char width for 40pt font is roughly 20px
                        chars_per_line = max(10, int(box_w / 22)) 
                        lines = textwrap.wrap(step_text, width=chars_per_line)
                        
                        # Calculate total text block height for vertical centering
                        total_h = 0
                        line_heights = []
                        for line in lines:
                            bbox = draw.textbbox((0,0), line, font=step_font)
                            lh = bbox[3] - bbox[1]
                            line_heights.append(lh)
                            total_h += lh
                        total_h += (len(lines) - 1) * 15 # padding between lines
                        
                        curr_y = y1 + (box_h - total_h) // 2
                        
                        # Draw centered text
                        for idx, line in enumerate(lines):
                            bbox = draw.textbbox((0,0), line, font=step_font)
                            line_w = bbox[2] - bbox[0]
                            draw.text((x1 + (box_w - line_w) // 2, curr_y), line, font=step_font, fill=(255, 255, 255))
                            curr_y += line_heights[idx] + 15
                            
                        # Draw arrow pointing to next box
                        if i < num_steps - 1:
                            arrow_x_start = x2 + 5
                            arrow_x_end = arrow_x_start + spacing - 10
                            arrow_y = y1 + box_h // 2
                            
                            draw.line([(arrow_x_start, arrow_y), (arrow_x_end, arrow_y)], fill=(15, 23, 42), width=8)
                            draw.polygon([(arrow_x_end, arrow_y), (arrow_x_end - 15, arrow_y - 15), (arrow_x_end - 15, arrow_y + 15)], fill=(15, 23, 42))
                            
            else:
                items = data.get("items", [])
                if not items:
                    # Fallback string
                    raw_text = scene.slide_body
                    if "•" in raw_text:
                        items = [{"title": "", "content": t.strip()} for t in raw_text.split("•") if t.strip()]
                    else:
                        items = [{"title": "", "content": raw_text}]
                
                start_y = 300
                for item in items:
                    t = item.get("title", "")
                    c = item.get("content", "")
                    text = f"{t}: {c}" if t and c else f"{t}{c}"
                    
                    start_y = draw_wrapped_text(draw, text, body_font, 150, start_y, 65, body_color, align="left", bullet=True)
                    start_y += 30 # space between bullets
                
        img.save(output_path)



class VideoComposer:
    def __init__(self, config: dict):
        self.config = config
        self.w, self.h = config['video']['resolution']
        self.fps = config['video']['fps']

    def compose(self, audio_path: str, slides_dir: str, scenes: List[Scene], output_path: str):
        narration_audio = AudioFileClip(audio_path)
        total_duration = narration_audio.duration
        
        total_words = sum(len(s.spoken_text.split()) for s in scenes)
        words_per_sec = total_words / total_duration if total_duration > 0 else 1.0
        
        visual_clips = []
        current_time = 0.0
        
        logo_path = self.config['branding'].get('logo_path', '')
        logo_clip = None
        if logo_path and os.path.exists(logo_path):
            logo_clip = ImageClip(logo_path).resize(height=60).set_position(('right', 'top')).margin(right=50, top=50, opacity=0).set_opacity(0.6).set_duration(total_duration)

        for i, scene in enumerate(scenes):
            scene_words = len(scene.spoken_text.split())
            scene_duration = max(2.0, scene_words / words_per_sec)
            
            if i == len(scenes) - 1:
                scene_duration = max(2.0, total_duration - current_time)
                
            slide_img_path = os.path.join(slides_dir, f"slide_{i}.png")
            slide_clip = ImageClip(slide_img_path).set_start(current_time).set_duration(scene_duration)
            
            if i > 0:
                slide_clip = slide_clip.crossfadein(0.5)
                
            visual_clips.append(slide_clip)
            current_time += scene_duration
            
        if logo_clip:
            visual_clips.append(logo_clip)
            
        video = CompositeVideoClip(visual_clips, size=(self.w, self.h))
        video = video.set_audio(narration_audio)
        
        music_folder = self.config['music'].get('folder', '')
        if music_folder and os.path.exists(music_folder):
            music_files = glob.glob(os.path.join(music_folder, '*.mp3'))
            if music_files:
                bgm_path = random.choice(music_files)
                # Reduce the background music volume by half to make it less loud but still audible
                bgm_volume = self.config['music'].get('volume', 0.1) * 0.5
                bgm_clip = AudioFileClip(bgm_path).fx(afx.volumex, bgm_volume)
                
                if bgm_clip.duration < total_duration:
                    bgm_clip = afx.audio_loop(bgm_clip, duration=total_duration)
                else:
                    bgm_clip = bgm_clip.subclip(0, total_duration)
                    
                bgm_clip = bgm_clip.fx(afx.audio_fadeout, self.config['music']['fade_duration'])
                final_audio = CompositeAudioClip([video.audio, bgm_clip])
                video = video.set_audio(final_audio)

        print(f"Rendering video to {output_path} (Duration: {total_duration:.2f}s)")
        video.write_videofile(
            output_path, 
            fps=self.fps,
            codec='libx264',
            audio_codec='aac',
            threads=4,
            logger='bar'
        )
        video.close()
        narration_audio.close()

def create_poster(title, description, top_label, output_path, is_centered=False):
    """Generates a purely programmatic vector-style poster using Pillow."""
    width, height = 1080, 1920
    
    # Base background (Dark Blue)
    img = Image.new('RGB', (width, height), color=(12, 61, 109))
    draw = ImageDraw.Draw(img)
    
    # Draw geometric accents to make it look professional and designed
    # Accent 1: Top right vivid blue circle
    draw.ellipse([width - 500, -200, width + 300, 600], fill=(12, 136, 235))
    
    # Accent 2: Bottom left dark slate circle
    draw.ellipse([-400, height - 700, 400, height + 100], fill=(15, 23, 42))
    
    # Accent 3: Small decorative circle
    draw.ellipse([width - 150, height - 400, width - 50, height - 300], fill=(255, 200, 50))

    # Load fonts
    font_path = branding().font_path or "/System/Library/Fonts/Supplemental/Arial.ttf"
    try:
        font_title = ImageFont.truetype(font_path, 85)
        font_desc = ImageFont.truetype(font_path, 55)
        font_top = ImageFont.truetype(font_path, 45)
        font_footer = ImageFont.truetype(font_path, 40)
    except Exception as e:
        print(f"Warning: Could not load custom font, using default. ({e})")
        font_title = ImageFont.load_default()
        font_desc = ImageFont.load_default()
        font_top = ImageFont.load_default()
        font_footer = ImageFont.load_default()

    px = 100
    y_text = 500 if is_centered else 300

    # 1. Draw Top Label
    if top_label:
        draw.text((px, y_text), top_label, fill=(255, 200, 50), font=font_top)
        y_text += 100

    # 2. Draw Title
    title_lines = textwrap.wrap(title, width=18)
    for line in title_lines:
        draw.text((px, y_text), line, fill=(255, 255, 255), font=font_title)
        y_text += 100

    # Draw separator line
    y_text += 40
    draw.line([px, y_text, px + 200, y_text], fill=(12, 136, 235), width=10)
    y_text += 80

    # 3. Draw Description
    if description:
        desc_lines = textwrap.wrap(description, width=28)
        for line in desc_lines:
            draw.text((px, y_text), line, fill=(220, 220, 220), font=font_desc)
            y_text += 75

    # 4. Draw Footer
    draw.text((px, height - 150), "InmiBot.es | Trámites de Extranjería", fill=(255, 255, 255), font=font_footer)

    img.save(output_path)

